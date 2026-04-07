#!/usr/bin/env python3
"""Generate RoboHire partner sales-kit PowerPoint and PDF assets."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    Flowable,
    KeepTogether,
    ListFlowable,
    ListItem,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)


ROOT = Path(__file__).resolve().parents[1]
PPTX_PATH = ROOT / "output" / "pptx" / "RoboHire_Sales_Kit.pptx"
PDF_PATH = ROOT / "output" / "pdf" / "RoboHire_Sales_Kit.pdf"

NAVY = RGBColor(15, 23, 42)
SLATE = RGBColor(71, 85, 105)
MUTED = RGBColor(100, 116, 139)
BLUE = RGBColor(37, 99, 235)
CYAN = RGBColor(6, 182, 212)
INDIGO = RGBColor(79, 70, 229)
GREEN = RGBColor(16, 185, 129)
AMBER = RGBColor(245, 158, 11)
WHITE = RGBColor(255, 255, 255)
BG = RGBColor(248, 250, 252)
BORDER = RGBColor(226, 232, 240)

HEX = {
    "navy": "#0f172a",
    "slate": "#475569",
    "muted": "#64748b",
    "blue": "#2563eb",
    "cyan": "#06b6d4",
    "indigo": "#4f46e5",
    "green": "#10b981",
    "amber": "#f59e0b",
    "bg": "#f8fafc",
    "border": "#e2e8f0",
    "light_blue": "#eff6ff",
    "light_cyan": "#ecfeff",
}


def ensure_dirs() -> None:
    PPTX_PATH.parent.mkdir(parents=True, exist_ok=True)
    PDF_PATH.parent.mkdir(parents=True, exist_ok=True)


def add_slide(prs: Presentation, title: str, subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Decorative SaaS-style blobs.
    blob = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(9.2),
        Inches(-0.55),
        Inches(3.4),
        Inches(3.4),
    )
    blob.fill.solid()
    blob.fill.fore_color.rgb = RGBColor(219, 234, 254)
    blob.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(0.08),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = BLUE
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.62), Inches(0.45), Inches(8.4), Inches(0.5))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Aptos Display"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.65), Inches(0.98), Inches(7.8), Inches(0.35))
        tf = sub.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.name = "Aptos"
        p.font.size = Pt(10.5)
        p.font.color.rgb = MUTED

    footer = slide.shapes.add_textbox(Inches(0.65), Inches(7.05), Inches(5), Inches(0.2))
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "RoboHire Partner Sales Kit"
    p.font.name = "Aptos"
    p.font.size = Pt(8.5)
    p.font.color.rgb = MUTED

    return slide


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = SLATE,
    align=PP_ALIGN.LEFT,
    line_spacing: float | None = None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    p.font.name = "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def add_pill(slide, x: float, y: float, w: float, text: str, color: RGBColor = BLUE):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(0.34),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(239, 246, 255)
    shape.line.color.rgb = RGBColor(191, 219, 254)
    shape.line.width = Pt(0.75)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Aptos"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = color
    return shape


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    body: str,
    accent: RGBColor = BLUE,
    title_size: int = 14,
    body_size: int = 10,
):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)

    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(0.08),
        Inches(h),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    add_text(slide, x + 0.2, y + 0.18, w - 0.35, 0.28, title, title_size, True, NAVY)
    add_text(slide, x + 0.2, y + 0.57, w - 0.35, h - 0.72, body, body_size, False, SLATE)
    return shape


def add_metric(slide, x: float, y: float, w: float, value: str, label: str, accent: RGBColor):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(0.9),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(248, 250, 252)
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)
    add_text(slide, x + 0.16, y + 0.15, w - 0.32, 0.32, value, 20, True, accent, PP_ALIGN.CENTER)
    add_text(slide, x + 0.16, y + 0.53, w - 0.32, 0.22, label, 8.5, True, MUTED, PP_ALIGN.CENTER)


def add_bullets(slide, x: float, y: float, w: float, h: float, bullets: list[str], size: int = 12):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = SLATE
        p.space_after = Pt(6)
        p._p.get_or_add_pPr().set("marL", "230000")
        p._p.get_or_add_pPr().set("indent", "-130000")
    return box


def add_arrow(slide, x: float, y: float, w: float, color: RGBColor = BLUE):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(0.22),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def build_pptx() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide = add_slide(prs, "")
    add_pill(slide, 0.72, 0.72, 1.85, "AI Recruiting Agents", BLUE)
    add_text(slide, 0.72, 1.35, 6.1, 1.65, "AI Screens.\nAI Interviews.\nYou Hire the Best.", 34, True, NAVY)
    add_text(
        slide,
        0.76,
        3.35,
        5.75,
        0.82,
        "Partner-ready product introduction kit for startup founders, HR teams, and hiring partners.",
        15,
        False,
        SLATE,
    )
    add_metric(slide, 0.76, 4.52, 1.65, "80%", "first-funnel work", BLUE)
    add_metric(slide, 2.62, 4.52, 1.65, "24/7", "AI interviews", CYAN)
    add_metric(slide, 4.48, 4.52, 1.65, "7", "languages", INDIGO)
    add_text(slide, 0.76, 6.18, 5.4, 0.3, "Start free or book a demo at app.robohire.io", 12, True, BLUE)

    add_card(slide, 7.35, 1.05, 4.95, 1.0, "1. Clarify role", "AI turns vague needs into a structured hiring brief.", CYAN)
    add_arrow(slide, 9.55, 2.17, 0.8, CYAN)
    add_card(slide, 7.35, 2.62, 4.95, 1.0, "2. Screen resumes", "Rank candidates by fit, gaps, strengths, and reasoning.", BLUE)
    add_arrow(slide, 9.55, 3.74, 0.8, BLUE)
    add_card(slide, 7.35, 4.19, 4.95, 1.0, "3. Interview and evaluate", "Run AI interviews and produce structured scorecards.", INDIGO)

    # Slide 2
    slide = add_slide(prs, "Startup Hiring Is Too Slow", "The first mile is messy, repetitive, and expensive for lean teams.")
    add_metric(slide, 0.75, 1.45, 2.3, "42 days", "typical hiring cycle", BLUE)
    add_metric(slide, 3.25, 1.45, 2.3, "200+", "resumes per role", CYAN)
    add_metric(slide, 5.75, 1.45, 2.3, "0-1", "dedicated recruiters", INDIGO)
    add_metric(slide, 8.25, 1.45, 2.3, "2+ weeks", "calendar drag", AMBER)
    pains = [
        ("Fuzzy requirements", "Founders and HR spend days translating rough needs into real role criteria."),
        ("Resume overload", "High applicant volume hides the few people who are actually worth advancing."),
        ("Scheduling drag", "Candidates, interviewers, calendars, and time zones slow the process down."),
        ("Inconsistent interviews", "Different interviewers ask different questions, so comparisons are hard to defend."),
    ]
    x_positions = [0.8, 3.95, 7.1, 10.25]
    for x, (title, body), color in zip(x_positions, pains, [BLUE, CYAN, INDIGO, AMBER]):
        add_card(slide, x, 3.0, 2.55, 2.05, title, body, color, 13, 9.5)
    add_text(
        slide,
        0.8,
        5.75,
        10.8,
        0.55,
        "RoboHire positions itself where startups feel the most pain: before the human interview loop, when teams need fast signal and a clean shortlist.",
        15,
        True,
        NAVY,
    )

    # Slide 3
    slide = add_slide(prs, "RoboHire Runs The First 80% Of Hiring", "AI agents handle the repetitive work before final human decision-making.")
    add_text(slide, 0.78, 1.45, 5.7, 1.0, "From vague role need to decision-ready shortlist.", 28, True, NAVY)
    add_bullets(
        slide,
        0.9,
        2.75,
        5.6,
        2.4,
        [
            "Clarifies the role like a strong recruiter.",
            "Generates the JD and structured criteria.",
            "Screens resumes with contextual reasoning.",
            "Invites candidates automatically.",
            "Runs AI-led first-round interviews 24/7.",
            "Creates structured scorecards for finalist review.",
        ],
        13,
    )
    outcomes = [
        ("Speed", "Move from application pile to shortlist in days, not weeks."),
        ("Consistency", "Use one structured rubric across first-round evaluation."),
        ("Leverage", "Give a small team recruiting capacity without adding headcount."),
    ]
    for i, (title, body) in enumerate(outcomes):
        add_card(slide, 7.15, 1.45 + i * 1.55, 4.85, 1.14, title, body, [BLUE, CYAN, INDIGO][i], 16, 11)
    add_text(slide, 7.2, 6.35, 4.75, 0.35, "Human team still owns final interviews and hiring decisions.", 11.5, True, MUTED)

    # Slide 4
    slide = add_slide(prs, "Workflow: Role Brief To Shortlist", "Six steps, one AI-driven hiring flow.")
    steps = [
        ("01", "Clarify", "AI role intake"),
        ("02", "Create JD", "Publishable draft"),
        ("03", "Screen", "Ranked candidates"),
        ("04", "Invite", "Private links"),
        ("05", "Interview", "24/7 AI video"),
        ("06", "Scorecard", "Decision report"),
    ]
    for i, (num, title, body) in enumerate(steps):
        x = 0.7 + i * 2.08
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.62), Inches(1.72), Inches(3.25))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = BORDER
        shape.line.width = Pt(1)
        add_text(slide, x + 0.18, 1.88, 1.35, 0.32, num, 13, True, [BLUE, CYAN, INDIGO, BLUE, CYAN, INDIGO][i], PP_ALIGN.CENTER)
        add_text(slide, x + 0.18, 2.45, 1.35, 0.55, title, 15, True, NAVY, PP_ALIGN.CENTER)
        add_text(slide, x + 0.16, 3.22, 1.42, 0.72, body, 10, False, SLATE, PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_arrow(slide, x + 1.75, 3.03, 0.38, BLUE)
    add_text(
        slide,
        0.8,
        5.55,
        10.9,
        0.7,
        "Best first pilot: choose one active role, upload a real resume batch, invite the top-fit candidates, and compare RoboHire scorecards against the team's manual shortlist.",
        15,
        True,
        NAVY,
    )

    # Slide 5
    slide = add_slide(prs, "What The AI Agents Do", "Feature menu for founders, HR teams, and integrations.")
    features = [
        ("AI Recruiting Consultant", "Turns messy hiring requests into structured role briefs."),
        ("AI Resume Screening", "Matches resumes to the JD with reasoning, gaps, and ranking."),
        ("Automated Outreach", "Sends interview links and QR codes to qualified candidates."),
        ("AI Video Interviewer", "Runs first-round interviews with live follow-up questions."),
        ("Evaluation Reports", "Summarizes fit, strengths, concerns, recommendation, and risks."),
        ("API + ATS Workflows", "Supports APIs, webhooks, and ATS integrations for teams that need scale."),
    ]
    for idx, (title, body) in enumerate(features):
        col = idx % 3
        row = idx // 3
        add_card(slide, 0.78 + col * 4.12, 1.45 + row * 2.15, 3.45, 1.56, title, body, [BLUE, CYAN, INDIGO][col], 13.5, 10.2)
    add_pill(slide, 0.82, 6.05, 2.2, "7 interview languages", CYAN)
    add_pill(slide, 3.22, 6.05, 2.9, "Cheating analysis", INDIGO)
    add_pill(slide, 6.32, 6.05, 3.35, "Greenhouse, Lever, Ashby, BambooHR, Workable", BLUE)

    # Slide 6
    slide = add_slide(prs, "Why Founders And HR Teams Care", "Translate features into buyer outcomes.")
    personas = [
        ("Founder / CEO", "Stop being the part-time recruiter. Get a sharper role brief and a finalist shortlist faster.", BLUE),
        ("HR / People Lead", "Reduce repetitive screening and run consistent first-round evaluations across candidates.", CYAN),
        ("Hiring Manager", "Spend interview time only on candidates with clear evidence and scorecards.", INDIGO),
        ("Partner / RPO", "Give startup members and recruiting clients a scalable AI screening and interview layer.", GREEN),
    ]
    for i, (title, body, color) in enumerate(personas):
        add_card(slide, 0.86 + (i % 2) * 6.1, 1.5 + (i // 2) * 2.05, 5.25, 1.52, title, body, color, 16, 11.5)
    add_text(
        slide,
        1.0,
        6.1,
        10.7,
        0.45,
        "Sales angle: RoboHire is recruiting leverage, not another dashboard. It creates usable signal before the human interview loop.",
        15,
        True,
        NAVY,
        PP_ALIGN.CENTER,
    )

    # Slide 7
    slide = add_slide(prs, "Typical Startup Sprint", "150 applicants to 3-4 finalists in a focused workflow.")
    timeline = [
        ("Mon AM", "AI clarifies role and drafts JD"),
        ("Mon Lunch", "150 resumes ranked into top 15"),
        ("Mon PM", "AI sends candidate interview links"),
        ("Wed", "12 AI interviews become 3-4 finalist reviews"),
    ]
    for i, (label, body) in enumerate(timeline):
        x = 0.9 + i * 3.05
        add_card(slide, x, 1.6, 2.45, 2.05, label, body, [BLUE, CYAN, INDIGO, GREEN][i], 17, 11)
        if i < 3:
            add_arrow(slide, x + 2.47, 2.48, 0.5, BLUE)
    add_metric(slide, 1.2, 4.65, 2.2, "150", "applicants uploaded", BLUE)
    add_metric(slide, 3.85, 4.65, 2.2, "15", "ranked top fits", CYAN)
    add_metric(slide, 6.5, 4.65, 2.2, "12", "AI interviews", INDIGO)
    add_metric(slide, 9.15, 4.65, 2.2, "3-4", "finalists", GREEN)
    add_text(slide, 0.9, 6.15, 11.3, 0.55, "The point is not to remove the hiring team. It is to make every human interview higher signal.", 15, True, NAVY, PP_ALIGN.CENTER)

    # Slide 8
    slide = add_slide(prs, "Why RoboHire Feels Different", "Clear alternatives messaging for partners.")
    diffs = [
        ("Not an ATS", "ATS stores candidate data. RoboHire actively moves screening and first-round interviews forward.", BLUE),
        ("Context, not keywords", "Reads resumes and interview answers semantically instead of only matching surface terms.", CYAN),
        ("Consistent rubric", "Creates cleaner comparisons and more defensible first-round evaluation.", INDIGO),
        ("Startup-friendly", "Professional recruiting process without a full HR department or large agency budget.", GREEN),
    ]
    for i, (title, body, color) in enumerate(diffs):
        add_card(slide, 0.9 + (i % 2) * 5.95, 1.45 + (i // 2) * 1.95, 5.0, 1.42, title, body, color, 15, 10.5)
    add_text(slide, 1.05, 5.65, 11, 0.38, "Integrates into real workflows: APIs, webhooks, ATS integrations, reports, and candidate interview links.", 15, True, NAVY, PP_ALIGN.CENTER)
    add_pill(slide, 1.0, 6.27, 1.45, "API", BLUE)
    add_pill(slide, 2.65, 6.27, 1.7, "Webhooks", CYAN)
    add_pill(slide, 4.55, 6.27, 1.45, "ATS", INDIGO)
    add_pill(slide, 6.2, 6.27, 1.75, "Scorecards", GREEN)
    add_pill(slide, 8.15, 6.27, 2.2, "Interview links", BLUE)

    # Slide 9
    slide = add_slide(prs, "Pricing Snapshot", "Confirm commercial terms before signed quotes.")
    plans = [
        ("Starter", "$29/mo", "1 seat, 3 roles, 15 interviews, scoring and summaries.", BLUE),
        ("Growth", "$199/mo", "Unlimited seats, 120 interviews, 240 resume matches, API interview creation.", CYAN),
        ("Business", "$399/mo", "280 interviews, 500 resume matches, analytics, white-label reports, cheating analysis.", INDIGO),
        ("Custom", "Contact us", "Unlimited options, custom workflows, ATS integrations, dedicated support.", GREEN),
    ]
    for i, (name, price, desc, color) in enumerate(plans):
        x = 0.75 + i * 3.07
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.55), Inches(2.65), Inches(3.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = BORDER
        add_text(slide, x + 0.22, 1.87, 2.15, 0.35, name, 15, True, NAVY, PP_ALIGN.CENTER)
        add_text(slide, x + 0.22, 2.42, 2.15, 0.45, price, 22, True, color, PP_ALIGN.CENTER)
        add_text(slide, x + 0.22, 3.2, 2.15, 0.78, desc, 9.5, False, SLATE, PP_ALIGN.CENTER)
    add_text(slide, 0.85, 5.22, 11.7, 0.42, "Pay per use: $0.40 per resume/JD match and $2.00 per AI interview.", 14, True, NAVY, PP_ALIGN.CENTER)
    add_pill(slide, 3.9, 5.95, 1.9, "14-day free trial", GREEN)
    add_pill(slide, 6.0, 5.95, 2.25, "No credit card required", BLUE)

    # Slide 10
    slide = add_slide(prs, "Partner CTA", "Make the first customer action simple.")
    add_text(slide, 0.88, 1.5, 5.7, 0.9, "Best next step: pick one active role and run the first screening batch through RoboHire.", 26, True, NAVY)
    add_bullets(
        slide,
        0.98,
        2.85,
        5.7,
        2.0,
        [
            "Share the one-line pitch and demo link.",
            "Ask how many resumes and open roles they have now.",
            "Offer a 20-minute demo focused on one real role.",
            "Start with product UI, API workflow, or ATS integration path.",
        ],
        13,
    )
    add_card(slide, 7.25, 1.45, 4.75, 1.15, "Trial", "Start free and test one live role.", GREEN, 17, 12)
    add_card(slide, 7.25, 2.95, 4.75, 1.15, "Demo", "Book a 20-minute walkthrough from role brief to shortlist.", BLUE, 17, 12)
    add_card(slide, 7.25, 4.45, 4.75, 1.15, "Partner share", "Forward the email, newsletter blurb, or WhatsApp copy from the Markdown kit.", CYAN, 17, 12)
    add_text(slide, 7.3, 6.38, 4.7, 0.25, "app.robohire.io  |  robohire.io/demo", 12, True, BLUE, PP_ALIGN.CENTER)

    prs.save(PPTX_PATH)


class HR(Flowable):
    def __init__(self, width: float, color=colors.HexColor(HEX["border"]), thickness: float = 1):
        super().__init__()
        self.width = width
        self.color = color
        self.thickness = thickness
        self.height = 0.05 * inch

    def draw(self):
        self.canv.setStrokeColor(self.color)
        self.canv.setLineWidth(self.thickness)
        self.canv.line(0, 0, self.width, 0)


def pdf_styles():
    styles = getSampleStyleSheet()
    styles.add(
        ParagraphStyle(
            name="CoverTitle",
            parent=styles["Title"],
            fontName="Helvetica-Bold",
            fontSize=30,
            leading=34,
            textColor=colors.HexColor(HEX["navy"]),
            spaceAfter=12,
        )
    )
    styles.add(
        ParagraphStyle(
            name="DeckTitle",
            parent=styles["Heading1"],
            fontName="Helvetica-Bold",
            fontSize=21,
            leading=25,
            textColor=colors.HexColor(HEX["navy"]),
            spaceAfter=12,
        )
    )
    styles.add(
        ParagraphStyle(
            name="DeckH2",
            parent=styles["Heading2"],
            fontName="Helvetica-Bold",
            fontSize=13,
            leading=16,
            textColor=colors.HexColor(HEX["blue"]),
            spaceBefore=8,
            spaceAfter=5,
        )
    )
    styles.add(
        ParagraphStyle(
            name="Body",
            parent=styles["BodyText"],
            fontName="Helvetica",
            fontSize=9.4,
            leading=13,
            textColor=colors.HexColor(HEX["slate"]),
            spaceAfter=6,
        )
    )
    styles.add(
        ParagraphStyle(
            name="Small",
            parent=styles["BodyText"],
            fontName="Helvetica",
            fontSize=8.3,
            leading=11,
            textColor=colors.HexColor(HEX["muted"]),
        )
    )
    styles.add(
        ParagraphStyle(
            name="Center",
            parent=styles["BodyText"],
            fontName="Helvetica-Bold",
            fontSize=10,
            leading=13,
            alignment=TA_CENTER,
            textColor=colors.HexColor(HEX["navy"]),
        )
    )
    styles.add(
        ParagraphStyle(
            name="Metric",
            parent=styles["BodyText"],
            fontName="Helvetica-Bold",
            fontSize=16,
            leading=18,
            alignment=TA_CENTER,
            textColor=colors.HexColor(HEX["blue"]),
        )
    )
    return styles


def P(text: str, style: ParagraphStyle) -> Paragraph:
    return Paragraph(text, style)


def bullet_list(items: list[str], styles) -> ListFlowable:
    return ListFlowable(
        [ListItem(P(item, styles["Body"]), leftIndent=10) for item in items],
        bulletType="bullet",
        bulletFontName="Helvetica",
        bulletFontSize=8,
        start="circle",
        leftIndent=14,
        bulletOffsetY=1,
    )


def styled_table(data, col_widths, header: bool = True, bg: str = "#ffffff") -> Table:
    table = Table(data, colWidths=col_widths, hAlign="LEFT")
    style = TableStyle(
        [
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("BOX", (0, 0), (-1, -1), 0.6, colors.HexColor(HEX["border"])),
            ("INNERGRID", (0, 0), (-1, -1), 0.4, colors.HexColor(HEX["border"])),
            ("LEFTPADDING", (0, 0), (-1, -1), 8),
            ("RIGHTPADDING", (0, 0), (-1, -1), 8),
            ("TOPPADDING", (0, 0), (-1, -1), 7),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 7),
            ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor(bg)),
        ]
    )
    if header:
        style.add("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(HEX["light_blue"]))
        style.add("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor(HEX["navy"]))
        style.add("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold")
    table.setStyle(style)
    return table


def metric_table(styles) -> Table:
    data = [
        [P("80%", styles["Metric"]), P("24/7", styles["Metric"]), P("7", styles["Metric"]), P("14 days", styles["Metric"])],
        [
            P("first-funnel work", styles["Small"]),
            P("AI interviews", styles["Small"]),
            P("interview languages", styles["Small"]),
            P("free trial", styles["Small"]),
        ],
    ]
    table = Table(data, colWidths=[1.45 * inch] * 4)
    table.setStyle(
        TableStyle(
            [
                ("BOX", (0, 0), (-1, -1), 0.6, colors.HexColor(HEX["border"])),
                ("INNERGRID", (0, 0), (-1, -1), 0.4, colors.HexColor(HEX["border"])),
                ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor(HEX["bg"])),
                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("TOPPADDING", (0, 0), (-1, -1), 8),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
            ]
        )
    )
    return table


def header_footer(canvas, doc) -> None:
    canvas.saveState()
    width, height = letter
    canvas.setFillColor(colors.HexColor(HEX["blue"]))
    canvas.rect(0, height - 0.08 * inch, width, 0.08 * inch, stroke=0, fill=1)
    canvas.setFillColor(colors.HexColor(HEX["muted"]))
    canvas.setFont("Helvetica", 8)
    canvas.drawString(doc.leftMargin, 0.45 * inch, "RoboHire Partner Sales Kit")
    canvas.drawRightString(width - doc.rightMargin, 0.45 * inch, f"Page {doc.page}")
    canvas.restoreState()


def section_title(text: str, styles):
    return KeepTogether([P(text, styles["DeckTitle"]), HR(6.3 * inch), Spacer(1, 0.12 * inch)])


def build_pdf() -> None:
    styles = pdf_styles()
    doc = SimpleDocTemplate(
        str(PDF_PATH),
        pagesize=letter,
        rightMargin=0.72 * inch,
        leftMargin=0.72 * inch,
        topMargin=0.72 * inch,
        bottomMargin=0.68 * inch,
        title="RoboHire Partner Sales Kit",
        author="RoboHire",
    )

    story = []

    story.append(Spacer(1, 0.35 * inch))
    story.append(P("RoboHire Partner Sales Kit", styles["CoverTitle"]))
    story.append(P("AI Screens. AI Interviews. You Hire the Best.", styles["DeckH2"]))
    story.append(
        P(
            "Partner-ready product introduction material for startup founders, HR teams, hiring managers, startup programs, and recruiting partners.",
            styles["Body"],
        )
    )
    story.append(Spacer(1, 0.18 * inch))
    story.append(metric_table(styles))
    story.append(Spacer(1, 0.22 * inch))
    story.append(
        P(
            "<b>Executive one-liner:</b> RoboHire is an AI recruiting team for startups and lean HR teams. It automates role scoping, JD generation, resume screening, candidate outreach, AI video interviews, and structured evaluation reports, so teams spend time only on the best-fit finalists.",
            styles["Body"],
        )
    )
    story.append(P("<b>Primary CTA:</b> Start a 14-day free trial or book a demo.", styles["Body"]))
    story.append(P("<b>Links:</b> app.robohire.io | robohire.io/demo | app.robohire.io/docs", styles["Body"]))
    story.append(PageBreak())

    story.append(section_title("1. The Problem And The Pitch", styles))
    story.append(
        P(
            "Hiring is too slow for startups. Founders and HR teams lose days clarifying roles, reviewing hundreds of resumes, coordinating interviews, and trying to compare candidates fairly.",
            styles["Body"],
        )
    )
    story.append(
        P(
            "RoboHire runs the first 80% of hiring on autopilot: role scoping, job description drafting, contextual resume screening, candidate invitations, AI-led first-round interviews, and structured scorecards.",
            styles["Body"],
        )
    )
    story.append(P("Common pain points", styles["DeckH2"]))
    story.append(
        bullet_list(
            [
                "Role requirements are vague and require repeated clarification.",
                "One open role can receive 200+ resumes, but only a few are worth advancing.",
                "First-round scheduling creates delay, especially across time zones.",
                "Different interviewers ask different questions and score inconsistently.",
                "Startups often lack a dedicated recruiter or trained interview panel.",
            ],
            styles,
        )
    )
    story.append(Spacer(1, 0.12 * inch))
    buyer_data = [
        [P("<b>Buyer</b>", styles["Body"]), P("<b>Message</b>", styles["Body"])],
        [P("Founder / CEO", styles["Body"]), P("Stop becoming the part-time recruiter. Turn messy hiring needs into a shortlist of evaluated candidates.", styles["Body"])],
        [P("HR / People Lead", styles["Body"]), P("Reduce repetitive screening and use one consistent rubric for first-round evaluation.", styles["Body"])],
        [P("Hiring Manager", styles["Body"]), P("Spend team interview time only on candidates with evidence and scorecards.", styles["Body"])],
        [P("Partner / RPO", styles["Body"]), P("Give startup members and recruiting clients a scalable AI screening and interview layer.", styles["Body"])],
    ]
    story.append(styled_table(buyer_data, [1.45 * inch, 4.85 * inch]))
    story.append(PageBreak())

    story.append(section_title("2. How RoboHire Works", styles))
    workflow = [
        [P("<b>Stage</b>", styles["Body"]), P("<b>What RoboHire Does</b>", styles["Body"]), P("<b>Outcome</b>", styles["Body"])],
        [P("1. Clarify", styles["Body"]), P("AI recruiting consultant asks about scope, must-haves, seniority, compensation, and fit.", styles["Body"]), P("Structured brief in about 10 minutes.", styles["Body"])],
        [P("2. Create JD", styles["Body"]), P("AI drafts responsibilities, requirements, and nice-to-haves.", styles["Body"]), P("Faster job launch.", styles["Body"])],
        [P("3. Screen", styles["Body"]), P("AI ranks candidates by context, skills, gaps, and potential.", styles["Body"]), P("Shortlist in minutes.", styles["Body"])],
        [P("4. Invite", styles["Body"]), P("AI sends private interview links and QR codes.", styles["Body"]), P("Less scheduling back-and-forth.", styles["Body"])],
        [P("5. Interview", styles["Body"]), P("AI video interviewer runs 24/7 with follow-up questions and 7 languages.", styles["Body"]), P("Candidates self-serve first round.", styles["Body"])],
        [P("6. Evaluate", styles["Body"]), P("Scorecards cover fit, experience, strengths, concerns, recommendations, and cheating analysis.", styles["Body"]), P("Decision-ready finalist review.", styles["Body"])],
    ]
    story.append(styled_table(workflow, [0.9 * inch, 3.25 * inch, 2.15 * inch]))
    story.append(Spacer(1, 0.15 * inch))
    story.append(P("Feature menu", styles["DeckH2"]))
    story.append(
        bullet_list(
            [
                "AI Recruiting Consultant and AI JD Generator.",
                "AI Resume Screening with match reasoning and gap analysis.",
                "Automated interview invitations with private links and QR codes.",
                "AI Video Interviewer with multilingual first-round interviews.",
                "Evaluation Reports with structured recommendations and cheating analysis.",
                "APIs, webhooks, and ATS integrations including Greenhouse, Lever, Ashby, BambooHR, and Workable.",
            ],
            styles,
        )
    )
    story.append(PageBreak())

    story.append(section_title("3. Startup Sprint Story", styles))
    story.append(
        P(
            "Use this narrative when prospects ask what the product feels like in practice.",
            styles["Body"],
        )
    )
    sprint = [
        [P("<b>Moment</b>", styles["Body"]), P("<b>RoboHire Action</b>", styles["Body"])],
        [P("Monday morning", styles["Body"]), P("Founder tells the AI recruiting consultant they need a senior product manager. RoboHire clarifies scope, must-have skills, seniority, compensation, and team fit.", styles["Body"])],
        [P("Before lunch", styles["Body"]), P("Team uploads 150 resumes. RoboHire returns a ranked top 15 with reasoning, strengths, and gaps.", styles["Body"])],
        [P("Monday afternoon", styles["Body"]), P("RoboHire automatically invites the best-fit candidates to complete first-round AI interviews.", styles["Body"])],
        [P("Wednesday", styles["Body"]), P("Completed interviews become structured scorecards. The team meets only the final 3 to 4 candidates.", styles["Body"])],
    ]
    story.append(styled_table(sprint, [1.35 * inch, 4.95 * inch]))
    story.append(Spacer(1, 0.14 * inch))
    story.append(P("Differentiators", styles["DeckH2"]))
    differentiators = [
        [P("<b>Not an ATS</b><br/>Most recruiting software stores information. RoboHire moves the work forward.", styles["Body"]), P("<b>Context, not keywords</b><br/>It reads resumes and interview answers semantically.", styles["Body"])],
        [P("<b>Consistent rubric</b><br/>Candidates can be compared against a shared first-round standard.", styles["Body"]), P("<b>Startup-friendly</b><br/>Run a professional process without a large recruiting department.", styles["Body"])],
    ]
    story.append(styled_table(differentiators, [3.1 * inch, 3.1 * inch], header=False, bg=HEX["bg"]))
    story.append(PageBreak())

    story.append(section_title("4. Pricing And Trial", styles))
    pricing = [
        [P("<b>Plan</b>", styles["Body"]), P("<b>Price</b>", styles["Body"]), P("<b>Best For</b>", styles["Body"]), P("<b>Highlights</b>", styles["Body"])],
        [P("Starter", styles["Body"]), P("$29/mo", styles["Body"]), P("Small teams", styles["Body"]), P("1 seat, 3 roles, 15 interviews, scoring and summaries.", styles["Body"])],
        [P("Growth", styles["Body"]), P("$199/mo", styles["Body"]), P("Scale mode", styles["Body"]), P("Unlimited seats, 120 interviews, 240 resume matches, API interview creation.", styles["Body"])],
        [P("Business", styles["Body"]), P("$399/mo", styles["Body"]), P("High-volume hiring", styles["Body"]), P("280 interviews, 500 resume matches, analytics, white-label reports, video playback, cheating analysis.", styles["Body"])],
        [P("Custom", styles["Body"]), P("Contact us", styles["Body"]), P("Large teams", styles["Body"]), P("Custom workflows, ATS integrations, unlimited options, dedicated manager and support.", styles["Body"])],
    ]
    story.append(styled_table(pricing, [0.9 * inch, 0.95 * inch, 1.25 * inch, 3.2 * inch]))
    story.append(Spacer(1, 0.15 * inch))
    story.append(P("<b>Pay per use:</b> $0.40 per resume/JD match and $2.00 per AI interview.", styles["Body"]))
    story.append(P("<b>Trial:</b> 14-day free trial. No credit card required. Start in minutes.", styles["Body"]))
    story.append(P("<b>Partner note:</b> Confirm commercial terms before sending a signed quote.", styles["Small"]))
    story.append(Spacer(1, 0.18 * inch))
    story.append(P("Recommended CTA blocks", styles["DeckH2"]))
    story.append(
        bullet_list(
            [
                "Start your free RoboHire trial and run one active role through AI screening and AI interviews.",
                "Book a 20-minute demo to see role brief to ranked shortlist.",
                "If your portfolio companies are hiring, share RoboHire as an AI recruiting accelerator for lean teams.",
            ],
            styles,
        )
    )
    story.append(PageBreak())

    story.append(section_title("5. Partner Copy And Guardrails", styles))
    story.append(P("Short email introduction", styles["DeckH2"]))
    story.append(
        styled_table(
            [
                [
                    P(
                        "<b>Subject:</b> AI hiring tool for faster startup recruiting<br/><br/>Hi {{first_name}},<br/><br/>I wanted to introduce RoboHire, an AI recruiting platform built for startups and lean HR teams.<br/><br/>RoboHire helps teams clarify roles, generate JDs, screen resumes, invite candidates, run AI-led first-round interviews, and produce structured candidate scorecards. The goal is to take the repetitive first 80% of hiring off the team, so founders and HR can focus on final interviews and decisions.<br/><br/>You can start with a free trial or book a quick demo here: {{demo_or_trial_link}}",
                        styles["Body"],
                    )
                ]
            ],
            [6.3 * inch],
            header=False,
            bg=HEX["bg"],
        )
    )
    story.append(Spacer(1, 0.15 * inch))
    story.append(P("Partner guardrails", styles["DeckH2"]))
    guardrails = [
        [P("<b>Use</b>", styles["Body"]), P("<b>Avoid</b>", styles["Body"])],
        [P("RoboHire helps automate the first 80% of hiring.", styles["Body"]), P("Guaranteed to hire in 3 days.", styles["Body"])],
        [P("Designed to move teams from resume overload to shortlist faster.", styles["Body"]), P("Fully removes human hiring decisions.", styles["Body"])],
        [P("AI-led first-round interviews with structured scorecards.", styles["Body"]), P("Legally eliminates bias or guarantees compliance.", styles["Body"])],
    ]
    story.append(styled_table(guardrails, [3.15 * inch, 3.15 * inch]))

    doc.build(story, onFirstPage=header_footer, onLaterPages=header_footer)


def main() -> None:
    ensure_dirs()
    build_pptx()
    build_pdf()
    print(PPTX_PATH)
    print(PDF_PATH)


if __name__ == "__main__":
    main()
